"""Build a deck from a plan and the stills captured on the same drive.

    build-plan-deck.py <plan.json> <stills-dir> <out.pptx>

The deck answers the same brief as the video, in the form you send ahead of a
conversation. It is NOT a second production run: the stills come from the recording pass
that already happened, so asking for both costs one extra pass rather than two drives —
which matters on a recording account whose project slots cannot be reclaimed.

Each scene becomes one slide: its caption is the argument, its still is the evidence.
That mapping is why the planner is told to caption the substance rather than the
mechanics — "Payments are held until the work is verified" is a slide title; "clicking
the payments tab" is not.
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

plan_path, stills_dir, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
plan = json.load(open(plan_path))

INK   = RGBColor(0x0A, 0x0A, 0x0A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
W, H  = Inches(13.333), Inches(7.5)

p = Presentation(); p.slide_width, p.slide_height = W, H


def blank():
    s = p.slides.add_slide(p.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PAPER
    return s


def text(s, txt, top, size, bold=False, colour=INK, left=Inches(0.9), width=None):
    tb = s.shapes.add_textbox(left, top, width or (W - Inches(1.8)), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True; tf.text = txt
    # Every paragraph and every run: `tf.text` splits on newlines, and styling only the
    # first paragraph leaves the rest at the default size — which is how the first
    # contractor deck came out with half its headlines looking like captions.
    for pr in tf.paragraphs:
        pr.alignment = PP_ALIGN.LEFT
        pr.line_spacing = 1.05
        for r in pr.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = colour; r.font.name = 'Arial'
    return tb


# Title slide, from the plan itself so the deck names the same thing the video does.
title = plan.get('title') or 'Groundwork by Jalla'
s = blank()
text(s, title, Inches(2.6), 40, True)
text(s, plan.get('subtitle') or 'Groundwork by Jalla', Inches(4.3), 18, False, MUTED)

IMG_W = Inches(7.0)
IMG_H = Emu(int(IMG_W * 810 / 1440))          # stills are 1440x810
img_top = Emu(int((H - IMG_H) / 2))

made = 0
for i, scene in enumerate(plan.get('scenes', [])):
    caption = scene.get('caption')
    still   = os.path.join(stills_dir, f'{i:02d}.png')
    # A slide needs both halves. A scene with no caption has no argument to make, and one
    # with no still has no evidence — either way it is better left out than padded.
    if not caption or not os.path.exists(still):
        continue
    s = blank()
    text(s, caption, Inches(2.9), 24, True, width=Inches(4.6))
    pic = s.shapes.add_picture(still, Inches(5.8), img_top, width=IMG_W)
    pic.line.color.rgb = RGBColor(0xE3, 0xE3, 0xE3); pic.line.width = Pt(0.75)
    made += 1

p.save(out_path)
print(f'wrote {out_path} — {made + 1} slides ({made} from scenes)', flush=True)
if made == 0:
    print('  !! no scene produced a slide — check captions and stills', file=sys.stderr)
