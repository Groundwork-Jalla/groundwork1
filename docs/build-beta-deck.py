"""Groundwork by Jalla — beta testing walkthrough deck.

Monochrome, matching the product: colour in Groundwork means STATE (active, held,
overdue), so a deck about the product should not spend colour on decoration either.
Hierarchy comes from weight, scale and space instead.
"""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK    = C(0x0A, 0x0A, 0x0A)
PAPER  = C(0xFA, 0xFA, 0xF8)
WHITE  = C(0xFF, 0xFF, 0xFF)
MID    = C(0x6B, 0x6B, 0x6B)
LIGHT  = C(0xD8, 0xD6, 0xD2)
FAINT  = C(0xF0, 0xEE, 0xEA)
ALERT  = C(0x9B, 0x22, 0x22)

W, H = In(13.333), In(7.5)
prs = Presentation(); prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

M = In(0.9)          # side margin
CW = W - 2*M         # content width

def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background()
    r.shadow.inherit = False
    return s

def box(s, x, y, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:    sh.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    else:    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def text(s, x, y, w, h, runs, size=16, color=INK, bold=False, align=PP_ALIGN.LEFT,
         spacing=1.15, anchor=MSO_ANCHOR.TOP, space_after=0):
    """runs: a string, or a list of (text, {overrides}) tuples."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    p.line_spacing = spacing; p.space_after = Pt(space_after)
    items = runs if isinstance(runs, list) else [(runs, {})]
    for t, ov in items:
        r = p.add_run(); r.text = t
        f = r.font
        f.name = 'Calibri'
        f.size = Pt(ov.get('size', size))
        f.bold = ov.get('bold', bold)
        f.color.rgb = ov.get('color', color)
    return tb

def bullet_list(s, x, y, w, items, size=15, gap=0.44, marker='—', color=INK, mcolor=MID):
    for i, it in enumerate(items):
        yy = y + In(gap * i)
        text(s, x, yy, In(0.32), In(0.4), marker, size=size, color=mcolor)
        text(s, x + In(0.36), yy, w - In(0.36), In(0.4), it, size=size, color=color)

def header(s, eyebrow, title, rule=True):
    if eyebrow:
        text(s, M, In(0.62), CW, In(0.3),
             eyebrow.upper(), size=11.5, color=MID, bold=True)
    text(s, M, In(0.98), CW, In(0.7), title, size=33, bold=True)
    if rule:
        box(s, M, In(1.72), In(1.5), Emu(22860), fill=INK)

def footer(s, n):
    text(s, M, H - In(0.62), In(6), In(0.3),
         'Groundwork by Jalla  ·  Beta testing walkthrough', size=9.5, color=MID)
    text(s, W - M - In(1), H - In(0.62), In(1), In(0.3),
         str(n), size=9.5, color=MID, align=PP_ALIGN.RIGHT)

N = [0]
def page(s):
    N[0] += 1; footer(s, N[0]); return s

# ── 1. Title ─────────────────────────────────────────────
s = slide(INK)
text(s, M, In(2.25), CW, In(0.4), 'GROUNDWORK BY JALLA', size=13, color=C(0x9A,0x9A,0x9A), bold=True)
text(s, M, In(2.75), CW, In(1.6),
     'Testing the platform,\nbeginning to end', size=48, bold=True, color=WHITE, spacing=1.05)
box(s, M, In(4.6), In(2.2), Emu(28575), fill=C(0x55,0x55,0x55))
text(s, M, In(4.95), In(9), In(0.9),
     'From creating an account to completing a build — ten flows, in the order a\nreal client would meet them.',
     size=17, color=C(0xB8,0xB8,0xB8), spacing=1.35)
text(s, M, H - In(1.05), CW, In(0.4),
     'Version 1.0   ·   16 August 2026   ·   tryjalla.com', size=12, color=C(0x77,0x77,0x77))

# ── 2. The journey ───────────────────────────────────────
s = page(slide())
header(s, 'Overview', 'The journey you are testing')
phases = [
    ('01', 'Account',    'Sign up, confirm\nemail, onboarding'),
    ('02', 'Create',     '11 steps, from\ncountry to budget'),
    ('03', 'Track',      '10 stages, 60\nsubstages, evidence'),
    ('04', 'Money',      'Costing, payments,\nmilestones'),
    ('05', 'Complete',   'Final handover,\narchive or delete'),
]
bw = (CW - In(0.5)*4) / 5
for i, (num, title, body) in enumerate(phases):
    x = M + (bw + In(0.5)) * i
    box(s, x, In(2.35), bw, In(2.7), fill=WHITE, line=LIGHT)
    text(s, x + In(0.3), In(2.62), bw, In(0.5), num, size=30, bold=True, color=LIGHT)
    text(s, x + In(0.3), In(3.28), bw - In(0.5), In(0.4), title, size=17, bold=True)
    text(s, x + In(0.3), In(3.75), bw - In(0.5), In(1.1), body, size=12.5, color=MID, spacing=1.3)
    if i < 4:
        text(s, x + bw + In(0.1), In(3.45), In(0.3), In(0.3), '→', size=17, color=LIGHT)
box(s, M, In(5.55), CW, In(0.95), fill=FAINT)
text(s, M + In(0.35), In(5.78), CW - In(0.7), In(0.6),
     [('Test in this order. ', {'bold': True}),
      ('Each phase depends on the one before it — you cannot track a build you have not created.', {})],
     size=14.5)

# ── 3. Ground rules ──────────────────────────────────────
s = page(slide())
header(s, 'Before you start', 'Three things to know')
rules = [
    ('It is free', 'Stay on the Self Verify plan throughout. Nothing here will charge you.'),
    ('No real payment details', 'Payment processing is not live. The Payments screens record what you say '
                                'you paid — they do not move money.'),
    ('Confusion is the finding', 'If a screen puzzles you, that is what we need to hear. Note what you '
                                 'expected instead. "I did not understand this number" beats a typo report.'),
]
for i, (t, b) in enumerate(rules):
    y = In(2.35) + In(1.35) * i
    box(s, M, y, Emu(34290), In(1.05), fill=INK)
    text(s, M + In(0.45), y, In(4.2), In(0.4), t, size=19, bold=True)
    text(s, M + In(0.45), y + In(0.42), CW - In(0.9), In(0.7), b, size=14, color=MID, spacing=1.3)

# ── 4. Reporting ─────────────────────────────────────────
s = page(slide())
header(s, 'Before you start', 'Reporting something')
text(s, M, In(2.3), In(5.6), In(0.4), 'Give us five things', size=19, bold=True)
bullet_list(s, M, In(2.95), In(5.6), [
    'Which flow and step  (e.g. "Flow 2, step 8")',
    'What you expected to happen',
    'What actually happened',
    'A screenshot if the screen looked wrong',
    'Browser, and phone or computer',
], size=14.5, gap=0.5)
box(s, M + In(6.4), In(2.3), CW - In(6.4), In(3.1), fill=WHITE, line=LIGHT)
text(s, M + In(6.75), In(2.62), CW - In(7.1), In(0.4), 'Red error messages', size=17, bold=True)
text(s, M + In(6.75), In(3.12), CW - In(7.1), In(2.0),
     'Copy the exact wording, including the code in brackets.\n\n'
     'That code tells us precisely where it broke — it is the difference between an hour '
     'of guessing and a five-minute fix.', size=14, color=MID, spacing=1.35)

# ── 5. Flow 1 ────────────────────────────────────────────
s = page(slide())
header(s, 'Phase 01  ·  Account', 'Sign up and first sign-in')
text(s, M, In(2.25), In(5.4), In(0.35), 'DO THIS', size=11, bold=True, color=MID)
bullet_list(s, M, In(2.75), In(5.4), [
    'Go to tryjalla.com',
    'Click Join for Free',
    'Email + password, or Continue with Google',
    'Click the confirmation link in your inbox',
    'Onboarding — pick your country',
    'You land on the dashboard',
], size=14.5, gap=0.46, marker='·')
x2 = M + In(6.2)
text(s, x2, In(2.25), CW - In(6.2), In(0.35), 'CHECK', size=11, bold=True, color=MID)
bullet_list(s, x2, In(2.75), CW - In(6.2), [
    'Does the email arrive in a minute? Spam?',
    'French-speaking country — does it offer French?',
    'Is an empty dashboard understandable?',
    'Log out and back in',
    'Forgot password — does the reset work?',
], size=14.5, gap=0.46, marker='·')
text(s, M, In(6.15), CW, In(0.4), 'Expected time: 3 minutes', size=13, color=MID)

# ── 6. Flow 2 — the 11 steps ─────────────────────────────
s = page(slide())
header(s, 'Phase 02  ·  Create', 'The wizard, step by step')
steps = [
    ('1','Country'), ('2','Project type'), ('3','Building type'), ('4','Floors'),
    ('5','Rooms'), ('6','Staff quarters'), ('7','Roof'), ('8','Details'),
    ('9','Summary'), ('10','Plan'), ('11','Confirm budget'),
]
cols, cw2 = 6, (CW - In(0.28)*5) / 6
for i, (n, t) in enumerate(steps):
    r_, c_ = divmod(i, cols)
    x = M + (cw2 + In(0.28)) * c_
    y = In(2.35) + In(1.22) * r_
    key = t in ('Details', 'Summary', 'Confirm budget')
    box(s, x, y, cw2, In(1.02), fill=INK if key else WHITE, line=None if key else LIGHT)
    text(s, x + In(0.22), y + In(0.13), cw2, In(0.3), n, size=12, bold=True,
         color=C(0x88,0x88,0x88) if key else LIGHT)
    # 12pt, not 13: 'Confirm budget' is the longest label and wrapped out of its card.
    text(s, x + In(0.22), y + In(0.44), cw2 - In(0.3), In(0.5), t, size=12, bold=True,
         color=WHITE if key else INK, spacing=1.1)
box(s, M, In(5.05), CW, In(1.3), fill=FAINT)
text(s, M + In(0.35), In(5.3), CW - In(0.7), In(1.0),
     [('The three filled steps are where we need you most. ', {'bold': True}),
      ('Step 8 sets the floor area that drives every figure. Step 9 shows the cost breakdown. '
       'Step 11 is where you accept or override the budget.', {})],
     size=14.5, color=INK, spacing=1.35)
text(s, M, In(6.5), CW, In(0.3), 'Expected time: 10 minutes', size=13, color=MID)

# ── 7. The drawing ───────────────────────────────────────
s = page(slide())
header(s, 'Phase 02  ·  Create', 'Watch the drawing build itself')
text(s, M, In(2.3), In(6.0), In(1.3),
     'As you answer, a technical drawing assembles on the right — floors stack, windows '
     'appear per room, the roof changes shape.', size=16.5, spacing=1.4)
# 3.85, not 3.5: at 16.5pt the paragraph above runs to three lines and the question
# was landing on top of its last line.
text(s, M, In(3.85), In(6.0), In(0.9),
     'Does it match what you think\nyou are describing?', size=20, bold=True, spacing=1.25)
box(s, M + In(6.7), In(2.3), CW - In(6.7), In(2.9), fill=WHITE, line=LIGHT)
text(s, M + In(7.05), In(2.6), CW - In(7.4), In(0.4), 'Tell us if', size=15, bold=True)
bullet_list(s, M + In(7.05), In(3.2), CW - In(7.35), [
    'The window count is wrong for your rooms',
    'The floors are wrong',
    'The roof shape is not what you chose',
], size=13.5, gap=0.48, marker='·')
box(s, M, In(5.45), CW, In(0.95), fill=FAINT)
text(s, M + In(0.35), In(5.68), CW - In(0.7), In(0.6),
     'It is the only feedback that shows whether we understood your answers.',
     size=14.5, spacing=1.3)

# ── 8. Floor area ────────────────────────────────────────
s = page(slide())
header(s, 'Phase 02  ·  Step 8', 'The floor area is the number to scrutinise')
box(s, M, In(2.3), In(5.9), In(2.5), fill=INK)
text(s, M + In(0.5), In(2.62), In(5.0), In(0.4), 'BENCHMARK', size=11, bold=True, color=C(0x88,0x88,0x88))
text(s, M + In(0.5), In(3.05), In(5.0), In(0.9), '120 m²', size=44, bold=True, color=WHITE)
text(s, M + In(0.5), In(3.95), In(5.0), In(0.7),
     'per floor, for a 5-bedroom\n2-storey semi-detached', size=14, color=C(0xAA,0xAA,0xAA), spacing=1.3)
x2 = M + In(6.4)
# Each block is placed below the WRAPPED height of the one above it, not below its
# nominal start — the bold line runs to two lines here and used to sit under the grey.
text(s, x2, In(2.35), CW - In(6.4), In(1.2),
     'We suggest this from your room list. It drives the entire budget — more than '
     'anything else on the screen.', size=16, spacing=1.4)
text(s, x2, In(3.62), CW - In(6.4), In(0.8),
     'Note it is the area of ONE floor,\nnot the whole house.', size=16.5, bold=True, spacing=1.3)
text(s, x2, In(4.62), CW - In(6.4), In(1.1),
     'Is that clear from the screen? If you know what your house actually measures, does '
     'our suggestion look sensible?', size=14.5, color=MID, spacing=1.4)

# ── 9. Four cost lines ───────────────────────────────────
s = page(slide())
header(s, 'Phase 02  ·  Step 9', 'The four lines you are quoted')
lines = [
    ('Construction',     'The build itself — roughly 60% material, 40% labour'),
    ('Design',           'Architectural and structural drawings, charged on built area'),
    ('Professional Fee', 'Groundwork’s fee: stage supervision and verification'),
    ('Permit',           'Planning approval and building permit'),
]
for i, (t, b) in enumerate(lines):
    y = In(2.3) + In(0.82) * i
    box(s, M, y, CW, In(0.68), fill=WHITE, line=LIGHT)
    box(s, M, y, Emu(45720), In(0.68), fill=INK)
    text(s, M + In(0.4), y + In(0.16), In(3.0), In(0.4), t, size=16, bold=True)
    text(s, M + In(3.7), y + In(0.19), CW - In(4.1), In(0.4), b, size=13.5, color=MID)
box(s, M, In(5.75), CW, In(0.95), fill=FAINT, line=INK, lw=1.25)
text(s, M + In(0.35), In(5.98), CW - In(0.7), In(0.6),
     [('Please check the four add up to the total shown. ', {'bold': True}),
      ('They should, to the cent. If any two figures disagree anywhere in the app, that is a serious bug.', {})],
     size=14.5)

# ── 10. Confirm budget ───────────────────────────────────
s = page(slide())
header(s, 'Phase 02  ·  Step 11', 'Confirm the budget')
text(s, M, In(2.3), In(6.2), In(1.6),
     'Accept our estimate, or type your contractor’s figure instead.\n\n'
     'Every stage payment is recalculated from whatever you confirm.', size=17, spacing=1.45)
box(s, M + In(6.9), In(2.3), CW - In(6.9), In(2.3), fill=WHITE, line=LIGHT)
text(s, M + In(7.25), In(2.6), CW - In(7.6), In(0.4), 'Try this', size=16, bold=True)
text(s, M + In(7.25), In(3.1), CW - In(7.6), In(1.3),
     'Type a figure 10% higher than ours, then check the payment schedule matches your '
     'number rather than ours.', size=14, color=MID, spacing=1.35)
box(s, M, In(4.9), CW, In(1.4), fill=INK)
text(s, M + In(0.45), In(5.15), CW - In(0.9), In(0.9),
     [('Free plan limit: three active projects.  ', {'bold': True, 'color': WHITE}),
      ('To create a fourth, archive an old one first. Archived projects do not count '
       'toward the limit.', {'color': C(0xBB,0xBB,0xBB)})],
     size=15, spacing=1.35)

# ── 11. Ten stages ───────────────────────────────────────
s = page(slide())
header(s, 'Phase 03  ·  Track', 'Ten stages, in build order')
stages = ['Land Secured','Design Completed','Site Preparation','Foundation','Structure & Walls',
          'Roofing','Electrical & Plumbing','Finishing','Exterior Work','Final Handover']
cw3 = (CW - In(0.22)*4) / 5
for i, name in enumerate(stages):
    r_, c_ = divmod(i, 5)
    x = M + (cw3 + In(0.22)) * c_
    y = In(2.35) + In(1.5) * r_
    unpaid = name in ('Land Secured', 'Exterior Work')
    box(s, x, y, cw3, In(1.25), fill=WHITE, line=LIGHT)
    text(s, x + In(0.24), y + In(0.16), cw3, In(0.35), f'{i+1:02d}', size=15, bold=True, color=LIGHT)
    text(s, x + In(0.24), y + In(0.55), cw3 - In(0.4), In(0.6), name, size=13.5, bold=True, spacing=1.2)
    if unpaid:
        text(s, x + In(0.24), y + In(0.98), cw3 - In(0.4), In(0.25), '$0 — not charged', size=10.5, color=MID)
text(s, M, In(5.58), CW, In(0.35),
     '60 substages in total. Only the current stage is unlocked — you cannot skip ahead, by design.',
     size=14.5, color=MID)
box(s, M, In(6.05), CW, In(0.72), fill=FAINT)
text(s, M + In(0.35), In(6.24), CW - In(0.7), In(0.4),
     'Land and Exterior Work carry no money. Does a $0 stage read as deliberate, or as a bug?',
     size=14)

# ── 12. Completing a stage ───────────────────────────────
s = page(slide())
header(s, 'Phase 03  ·  Track', 'How a stage completes')
seq = [
    ('01','Open the stage','Expand it to see its substages'),
    ('02','Upload evidence','A photo or document per substage'),
    ('03','Mark complete','Each substage in turn'),
    ('04','Approve the stage','Once every substage is done'),
    ('05','Next unlocks','The following stage opens'),
]
bw2 = (CW - In(0.3)*4) / 5
for i, (n, t, b) in enumerate(seq):
    x = M + (bw2 + In(0.3)) * i
    box(s, x, In(2.35), bw2, In(2.35), fill=INK if i == 4 else WHITE, line=None if i == 4 else LIGHT)
    text(s, x + In(0.26), In(2.6), bw2, In(0.35), n, size=13, bold=True,
         color=C(0x88,0x88,0x88) if i == 4 else LIGHT)
    text(s, x + In(0.26), In(3.05), bw2 - In(0.45), In(0.75), t, size=15, bold=True,
         color=WHITE if i == 4 else INK, spacing=1.2)
    text(s, x + In(0.26), In(3.85), bw2 - In(0.45), In(0.7), b, size=12, spacing=1.3,
         color=C(0xAA,0xAA,0xAA) if i == 4 else MID)
text(s, M, In(5.2), In(5.6), In(0.35), 'CHECK', size=11, bold=True, color=MID)
bullet_list(s, M, In(5.65), CW, [
    'Is it obvious what evidence each substage wants?',
    'Does a large photo straight from a phone camera upload?',
    'Do the stage names match how a builder in your country would say them?',
], size=14, gap=0.4, marker='·')

# ── 13. Money ────────────────────────────────────────────
s = page(slide())
header(s, 'Phase 04  ·  Money', 'Costing and payments')
text(s, M, In(2.25), In(5.6), In(0.35), 'COSTING TAB', size=11, bold=True, color=MID)
bullet_list(s, M, In(2.72), In(5.6), [
    'Open "How is this calculated?"',
    'Expand Construction to see it across the stages',
    'Export the PDF',
], size=14.5, gap=0.46, marker='·')
x2 = M + In(6.2)
text(s, x2, In(2.25), CW - In(6.2), In(0.35), 'PAYMENTS TAB', size=11, bold=True, color=MID)
bullet_list(s, x2, In(2.72), CW - In(6.2), [
    'Record a payment against Stage 1',
    'Watch the paid total and bar update',
    'Mark it unpaid again',
], size=14.5, gap=0.46, marker='·')
box(s, M, In(4.35), CW, In(1.5), fill=INK)
text(s, M + In(0.45), In(4.6), CW - In(0.9), In(1.05),
     [('The most valuable check in the whole beta:  ', {'bold': True, 'color': WHITE}),
      ('the total on the Overview donut, the Costing tab, the Payments schedule and the exported PDF '
       'must all be identical. If any two differ, screenshot both and send them.', {'color': C(0xBB,0xBB,0xBB)})],
     size=15.5, spacing=1.4)
text(s, M, In(6.15), CW, In(0.5),
     'Design, permit and professional fees appear as their own payment lines — they are not tied to site work.',
     size=13.5, color=MID)

# ── 14. Contractors ──────────────────────────────────────
s = page(slide())
header(s, 'Phase 05  ·  Contractors', 'Invitations and the Bill of Quantities')
text(s, M, In(2.25), In(5.6), In(0.35), 'AS THE HOMEOWNER', size=11, bold=True, color=MID)
bullet_list(s, M, In(2.72), In(5.6), [
    'Invite a contractor by email',
    'One contractor per project on the free plan',
    'They see Stages and Messages — not your budget',
], size=14, gap=0.46, marker='·')
x2 = M + In(6.2)
text(s, x2, In(2.25), CW - In(6.2), In(0.35), 'AS THE CONTRACTOR', size=11, bold=True, color=MID)
bullet_list(s, x2, In(2.72), CW - In(6.2), [
    'Open Bill of Quantities, start a take-off',
    'Every line arrives already priced',
    'Change a rate — the total moves immediately',
    'Save draft, then submit to the client',
], size=14, gap=0.46, marker='·')
box(s, M, In(4.75), CW, In(1.65), fill=FAINT)
text(s, M + In(0.4), In(4.98), CW - In(0.8), In(1.2),
     [('Contractors, this is the question we need answered:  ', {'bold': True}),
      ('do our item numbers match how you actually write a BQ? We use 204 footings, 305 blockwork, '
       '503 roof sheet, 801–810 plumbing. Lines marked "est." are rates we inferred rather than '
       'measured — tell us which are wrong.', {})],
     size=14.5, spacing=1.4)

# ── 15. Finish ───────────────────────────────────────────
s = page(slide())
header(s, 'Phase 05  ·  Complete', 'Finishing a project')
cards = [
    ('Final handover', 'Complete stage 10. The build is done and the project reads as complete.', False),
    ('Archive',        'Hides it and frees a plan slot. Everything is kept. Reversible at any time.', False),
    ('Delete',         'Permanent. Removes the project, stages, documents, payments and messages.', True),
]
bw3 = (CW - In(0.4)*2) / 3
for i, (t, b, danger) in enumerate(cards):
    x = M + (bw3 + In(0.4)) * i
    box(s, x, In(2.35), bw3, In(2.4), fill=WHITE, line=ALERT if danger else LIGHT, lw=1.5 if danger else 1.0)
    text(s, x + In(0.35), In(2.68), bw3 - In(0.6), In(0.45), t, size=19, bold=True,
         color=ALERT if danger else INK)
    text(s, x + In(0.35), In(3.28), bw3 - In(0.6), In(1.2), b, size=13.5, color=MID, spacing=1.35)
text(s, M, In(5.15), In(5.6), In(0.35), 'CHECK', size=11, bold=True, color=MID)
bullet_list(s, M, In(5.6), CW, [
    'Is the difference between archive and delete clear before you click?',
    'Does archiving free a slot, so you can create another project?',
    'Does the delete warning tell you enough to decide?',
], size=14, gap=0.4, marker='·')

# ── 16. Known limits ─────────────────────────────────────
s = page(slide())
header(s, 'Please note', 'Known limits — no need to report these')
limits = [
    ('Payments',          'Recording only. No money moves. Stripe and mobile money are not connected.'),
    ('Cost accuracy',     'Calibrated on four real Cameroonian bills of quantities. Cameroon is most accurate; '
                          'other countries are indexed estimates.'),
    ('Staff quarters',    'We ask the question but do not price it yet — awaiting a bill of quantities.'),
    ('Some BQ rates',     'Lines marked "est." have no source document behind them yet.'),
    ('Nigeria',           'Rates are unverified. We have no Nigerian bill of quantities.'),
    ('Jalla Management',  'Budget is set by our team after creation, so the project opens showing a banner.'),
]
for i, (t, b) in enumerate(limits):
    y = In(2.3) + In(0.72) * i
    box(s, M, y, CW, In(0.6), fill=WHITE if i % 2 == 0 else PAPER, line=LIGHT)
    text(s, M + In(0.35), y + In(0.14), In(2.9), In(0.35), t, size=14, bold=True)
    text(s, M + In(3.4), y + In(0.16), CW - In(3.75), In(0.35), b, size=12.5, color=MID)

# ── 17. Close ────────────────────────────────────────────
s = page(slide(INK))
text(s, M, In(2.3), CW, In(0.4), 'THANK YOU', size=13, bold=True, color=C(0x9A,0x9A,0x9A))
text(s, M, In(2.8), In(9.5), In(1.4),
     'Tell us what confused you.', size=42, bold=True, color=WHITE, spacing=1.1)
box(s, M, In(4.35), In(2.2), Emu(28575), fill=C(0x55,0x55,0x55))
text(s, M, In(4.7), In(9.5), In(1.4),
     'Every question in this deck exists because we genuinely do not know the answer.\n'
     'A confused tester is more useful to us than a polite one.',
     size=17, color=C(0xB8,0xB8,0xB8), spacing=1.4)
text(s, M, H - In(1.05), CW, In(0.4),
     'contact@tryjalla.com   ·   tryjalla.com', size=13, color=C(0x88,0x88,0x88))

prs.save('/home/favour-nwachukwu/Desktop/Jalla/groundwork1/docs/Groundwork-Beta-Walkthrough.pptx')
print('saved', N[0] + 2, 'slides')
