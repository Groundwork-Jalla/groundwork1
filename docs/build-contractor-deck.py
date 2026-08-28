"""Request b87fab1e, deck cut — "What is Groundwork by Jalla", for contractors.

The same argument as the video, in a form that can be sent ahead of a conversation or
left with someone. Built from live screenshots of the running app, not mockups.

Structure follows the brief's goal — value first, mechanics second, the ask last:
  1  the problem, in their words        4-6  where the work comes from
  2  what Groundwork is                 7    how they get paid
  3  what it changes                    8    how to apply
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SHOTS = sys.argv[1]
OUT   = os.path.abspath('docs/Groundwork-For-Contractors.pptx')

INK   = RGBColor(0x0A, 0x0A, 0x0A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
W, H  = Inches(13.333), Inches(7.5)

p = Presentation(); p.slide_width, p.slide_height = W, H

def blank():
    s = p.slides.add_slide(p.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = PAPER
    return s

def text(s, txt, top, size, bold=False, colour=INK, align=PP_ALIGN.LEFT,
         left=Inches(0.9), width=None):
    tb = s.shapes.add_textbox(left, top, width or (W - Inches(1.8)), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True; tf.text = txt
    # EVERY paragraph, every run. `tf.text = "a\nb"` splits on the newline into two
    # paragraphs, and styling only paragraphs[0] left the second line of every two-line
    # headline at the default 18pt regular — which is exactly how the first render came
    # out, with "The system around you is broken." looking like a caption.
    for pr in tf.paragraphs:
        pr.alignment = align
        pr.line_spacing = 1.05
        for r in pr.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = colour; r.font.name = 'Arial'
    return tb

def statement(title, sub=None, kicker=None):
    """A full-bleed argument slide. No image — these are the beats that carry the pitch."""
    s = blank()
    if kicker:
        text(s, kicker.upper(), Inches(1.4), 12, True, MUTED)
    text(s, title, Inches(2.1), 40, True)
    if sub:
        text(s, sub, Inches(3.9), 18, False, MUTED)
    return s

def shot(title, img, caption=None):
    """Screenshot right, argument left. The picture is evidence for the sentence.

    Both columns are optically centred on the same axis rather than hung from the top —
    the first render left roughly two inches of dead space under every screenshot, which
    reads as an unfinished slide rather than as breathing room.
    """
    s = blank()
    IMG_W = Inches(7.0)
    IMG_H = Emu(int(IMG_W * 810 / 1440))          # source is 1440x810
    img_top = Emu(int((H - IMG_H) / 2))
    text(s, title, Inches(2.35), 26, True, width=Inches(4.6))
    if caption:
        text(s, caption, Inches(3.45), 13, False, MUTED, width=Inches(4.6))
    path = os.path.join(SHOTS, img)
    if os.path.exists(path):
        pic = s.shapes.add_picture(path, Inches(5.8), img_top, width=IMG_W)
        pic.line.color.rgb = RGBColor(0xE3, 0xE3, 0xE3); pic.line.width = Pt(0.75)
    else:
        print(f'  !! missing {img}', file=sys.stderr)
    return s

# ── 1-3  the argument ────────────────────────────────────
statement("You're great at what you do.\nThe system around you is broken.",
          "Unclear scope. Payments 'handled later'. Projects that stall because nobody "
          "is coordinating the sequence.", kicker="For contractors")
statement("Groundwork is not a job board.",
          "It is controlled infrastructure for executing diaspora building projects "
          "properly — with the right professionals, in the right sequence, with the "
          "right safeguards.", kicker="What it is")
statement("Scope is agreed before you start.\nPayment is held before you build.",
          "The client funds each stage up front. Jalla verifies the work. The money is "
          "released against evidence, not against a phone call.", kicker="What changes")

# ── 4-6  where the work comes from ───────────────────────
shot("A client scopes the build before you are called",
     "05_wizard.png",
     "Eleven steps: location, building type, floors, rooms, roof, finish. The budget is "
     "measured from quantities at your city's rates — not a guess per square metre.")
shot("Every project arrives already costed",
     "06_project.png",
     "Construction, design, professional and permit, itemised. You are quoting against a "
     "figure the client has already accepted.")
shot("Ten stages, each with its own evidence",
     "07_stages.png",
     "Land, design, site prep, foundation, structure, roofing, services, finishing, "
     "exterior, handover. Everyone can see where the build actually is.")

# ── 7  the part that matters most ────────────────────────
shot("The money is already there",
     "08_payments.png",
     "Funds sit in escrow against the stage schedule. Complete the stage, submit the "
     "evidence, the payment is released. No chasing.")

# ── 8  the ask ───────────────────────────────────────────
shot("First in gets the best position",
     "03_apply.png",
     "We are onboarding a limited number of partners per trade, per region. The "
     "application takes about three minutes: your details, your credentials, and three "
     "past projects with references.")
statement("Apply to be a founding contractor",
          "tryjalla.com/contractor-apply", kicker="How to join")

p.save(OUT)
print('wrote', OUT, os.path.getsize(OUT), 'bytes,', len(p.slides.__iter__.__self__._sldIdLst), 'slides')
